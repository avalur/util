import pandas as pd
from pptx import Presentation

# read data
data = pd.read_csv('ismc-export-2020-06-30-fin.csv', sep=';')
data.head()

names = data[data['event'] == 'Вебинар по проверке кода'][['Name', 'Email']]
names.to_csv('participants_email.csv', index=False)

from pptx import Presentation

# participants cert
prs = Presentation('./templates/participant_ru.pptx')
first_slide = prs.slides[0]
shape = first_slide.shapes[0]
first_paragraph = shape.text_frame.paragraphs[0]

for name in names['Name']:
    first_paragraph.runs[0].text = name
    print(name)
    prs.save(f'./participant_pptx/{name}.pptx')

# speakers cert
prs = Presentation('./templates/speaker_ru.pptx')
first_slide = prs.slides[0]
shape = first_slide.shapes[1]
first_paragraph = shape.text_frame.paragraphs[0]

for name in ('Виталий Брагилевский', 'Алексей Толстиков', 'Артем Бурылов', 'Екатерина Лебедева', 'Сергей Жеревчук', 'Александр Паволоцкий', 'Илария Белова', 'Татьяна Васильева'):
    first_paragraph.runs[0].text = name
    print(name)
    prs.save(f'./speaker_pptx/{name}.pptx')

# make certs
from os import listdir

# import comtypes.client
import win32com.client

def PPTtoPDF(inputFileName, outputFileName, formatType=32):
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()
    
# %cd ./participant_pptx
# %cd ./speaker_pptx

# sometimes it did not work, but finally Ok
path_to_pptx = '.'
for pptx in listdir(path_to_pptx):
    inputFileName = pptx
    outputFileName = pptx.replace('pptx', 'pdf')
    print(inputFileName, outputFileName)
    PPTtoPDF(inputFileName, outputFileName)

# rename to pdf

import os

for f in listdir('.'):
    if f.endswith('pdf'):
        os.rename(f, f'../participant_pdf/{f}')
#         os.rename(f, f'../speaker_pdf/{f}')

